"""Generate deck.pptx — 10-slide design deck for the Q&A assistant."""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

DARK = RGBColor(0x0E, 0x1B, 0x2C)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
GREY = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(title: str, accent: bool = False):
    slide = prs.slides.add_slide(BLANK)
    if accent:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                    prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK
        bg.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18),
                                 prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = WHITE if accent else DARK
    return slide


def add_bullets(slide, items, top=1.6, left=0.8, width=11.8, size=19,
                color=GREY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                                  Inches(7.0 - top))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (level, text) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("•  " if level == 0 else "–  ") + text
        p.level = level
        p.font.size = Pt(size if level == 0 else size - 3)
        p.font.color.rgb = color
        p.space_after = Pt(10)
        p.line_spacing = line_spacing
    return tb


def box(slide, x, y, w, h, text, fill=TEAL, font_size=13, font_color=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                                 Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = DARK
    shp.line.width = Pt(0.75)
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = font_color
    return shp


def arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = GREY
    conn.line.width = Pt(2.25)
    le = conn.line._get_or_add_ln()
    le.append(  # arrowhead
        le.makeelement("{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd",
                       {"type": "arrow"}))


# ---- Slide 1: title -------------------------------------------------------
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()
tb = s.shapes.add_textbox(Inches(1), Inches(2.4), Inches(11.3), Inches(1.4))
p = tb.text_frame.paragraphs[0]
p.text = "Python Programming Q&A Assistant"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
tb2 = s.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.3), Inches(1.6))
tf = tb2.text_frame
p = tf.paragraphs[0]
p.text = "RAG over Stack Overflow's Python corpus · FastAPI · ChromaDB · Llama 3.1 (Groq)"
p.font.size = Pt(20)
p.font.color.rgb = TEAL
p2 = tf.add_paragraph()
p2.text = "AI Engineer Assessment — Round 1"
p2.font.size = Pt(16)
p2.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

# ---- Slide 2: what I built ------------------------------------------------
s = add_slide("What I built")
add_bullets(s, [
    (0, "A grounded Q&A API for data-science learners: ask any Python question, "
        "get an answer cited against real Stack Overflow answers."),
    (0, "POST /ask → answer + cited sources + grounding flag;  GET /health for probes."),
    (0, "30,000 highest-quality Q&A pairs indexed from the 987k-row Kaggle dataset "
        "(stackoverflow/pythonquestions)."),
    (0, "Fully tested: pytest suite with mocked LLM + 10 documented live queries "
        "covering edge cases."),
    (0, "Docker-ready; runs on free-tier hosts (no PyTorch — ONNX embeddings)."),
])

# ---- Slide 3: architecture (diagram) --------------------------------------
s = add_slide("Architecture")
# Online path
box(s, 0.7, 1.7, 1.6, 0.9, "Learner\nquestion", fill=GREY)
arrow(s, 2.3, 2.15, 3.0, 2.15)
box(s, 3.0, 1.55, 2.2, 1.2, "FastAPI /ask\nvalidate + cache")
arrow(s, 5.2, 2.15, 5.9, 2.15)
box(s, 5.9, 1.55, 2.2, 1.2, "Embed query\nMiniLM-L6 (ONNX)")
arrow(s, 8.1, 2.15, 8.8, 2.15)
box(s, 8.8, 1.55, 2.2, 1.2, "ChromaDB top-k\ncosine HNSW")
arrow(s, 9.9, 2.75, 9.9, 3.45)
box(s, 8.8, 3.45, 2.2, 1.2, "Relevance gate\n+ grounding prompt")
arrow(s, 8.8, 4.05, 8.1, 4.05)
box(s, 5.9, 3.45, 2.2, 1.2, "Groq Llama 3.1 8B\n(async, 30s timeout)")
arrow(s, 5.9, 4.05, 5.2, 4.05)
box(s, 3.0, 3.45, 2.2, 1.2, "Answer + [n] cites\n+ sources[]", fill=DARK)
# Offline path
tb = s.shapes.add_textbox(Inches(0.8), Inches(5.1), Inches(2), Inches(0.4))
p = tb.text_frame.paragraphs[0]
p.text = "Offline indexing:"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK
box(s, 0.7, 5.6, 2.3, 1.1, "Kaggle SO Python\n987k Q&A (parquet)", fill=GREY)
arrow(s, 3.0, 6.15, 3.6, 6.15)
box(s, 3.6, 5.6, 2.5, 1.1, "DuckDB: best answer\nper question", fill=GREY)
arrow(s, 6.1, 6.15, 6.7, 6.15)
box(s, 6.7, 5.6, 2.5, 1.1, "Top 30k by votes\nHTML → clean text", fill=GREY)
arrow(s, 9.2, 6.15, 9.8, 6.15)
box(s, 9.8, 5.6, 2.5, 1.1, "Embed title+question\n→ persistent index", fill=GREY)

# ---- Slide 4: design decisions 1 ------------------------------------------
s = add_slide("Design decisions — data & retrieval")
add_bullets(s, [
    (0, "Quality over volume: one best-voted answer per question, top 30k by score."),
    (1, "Community votes are a free reranking signal — every indexed answer is "
        "community-verified. Keeps index small, fast and trustworthy."),
    (0, "Asymmetric embedding: embed only title + question; store full Q&A for the LLM."),
    (1, "User queries resemble questions, not answers — and it sidesteps "
        "MiniLM's 256-token truncation."),
    (0, "ONNX MiniLM instead of PyTorch sentence-transformers."),
    (1, "~80 MB runtime vs ~2 GB; identical vectors; fits free-tier hosting."),
    (0, "DuckDB over raw parquet — no ETL infrastructure, one SQL window query."),
])

# ---- Slide 5: design decisions 2 ------------------------------------------
s = add_slide("Design decisions — generation & grounding")
add_bullets(s, [
    (0, "Llama 3.1 8B Instant on Groq: ~1 s end-to-end answers, generous free tier."),
    (0, "Grounding contract enforced in the system prompt:"),
    (1, "cite excerpts inline as [1], [2] · flag claims not from context · "
        "decline non-Python questions · modernise Python 2 answers (corpus is pre-2017)."),
    (0, "Relevance gate: if no hit clears the cosine threshold, the response is "
        "marked grounded: false instead of silently hallucinating."),
    (0, "Sources returned with every answer: SO link, vote count, similarity score — "
        "learners can verify everything."),
])

# ---- Slide 6: API design ---------------------------------------------------
s = add_slide("API design")
add_bullets(s, [
    (0, "POST /ask  {question, top_k?} → {answer, sources[], grounded, model, "
        "latency_ms, cached}"),
    (0, "GET /health → status + index size + model (readiness probe)"),
    (0, "Pydantic validation: question length 3–2000, top_k 1–10 → clean 422 errors."),
    (0, "Provider failures → 502 with detail; 30 s LLM timeout, never a hang."),
    (0, "In-process FIFO cache — repeated questions answered in 0 ms."),
    (0, "OpenAPI docs served at /docs."),
])

# ---- Slide 7: testing ------------------------------------------------------
s = add_slide("Testing")
add_bullets(s, [
    (0, "pytest, LLM mocked (runs without index or API key — CI-safe):"),
    (1, "health · happy path with citations · cache hit behaviour · empty / "
        "oversized / malformed input → 422 · provider failure → 502."),
    (0, "Retrieval integration tests against the real index:"),
    (1, "index populated · common questions retrieve relevant titles · "
        "off-topic queries score low."),
    (0, "10 live queries documented in TEST_RESULTS.md:"),
    (1, "basics, practical tasks, conceptual, pandas, vague phrasing, "
        "off-topic — each with observations on quality."),
])

# ---- Slide 8: observed behaviour -------------------------------------------
s = add_slide("Observed behaviour & edge cases")
add_bullets(s, [
    (0, "Strong on canonical questions (reverse a list, merge dicts, decorators): "
        "cited, runnable, correct."),
    (0, "Vague queries ('my code is slow how to make fast') still retrieve useful "
        "profiling threads — answer stays generic but grounded."),
    (0, "Off-topic ('capital of France') → politely declined, grounded: false."),
    (0, "Known limitation: corpus ends Oct 2016."),
    (1, "Prompt forces Python-3 modernisation; questions about newer features "
        "are answered from model knowledge and flagged as unsourced."),
])

# ---- Slide 9: scaling -------------------------------------------------------
s = add_slide("Scaling to 100+ concurrent users")
add_bullets(s, [
    (0, "Latency: LLM dominates (~1 s). Stream tokens via SSE for perceived speed; "
        "query embedding is ~5 ms on CPU."),
    (0, "Async: non-blocking Groq calls already; add uvicorn workers + load "
        "balancer — app is I/O-bound, LLM compute is Groq's problem."),
    (0, "Database: ChromaDB → managed vector DB (Qdrant / pgvector) so all "
        "replicas share one index; horizontal read scaling."),
    (0, "Caching: Redis exact-match + semantic cache — learner traffic repeats "
        "heavily, so hit rates are high."),
    (0, "Cost: 8B model ≈ $0.05 / 1M input tokens → well under $0.001 per answer; "
        "caching keeps 100 concurrent users at single-digit $ / day."),
])

# ---- Slide 10: summary ------------------------------------------------------
s = add_slide("Summary", accent=True)
add_bullets(s, [
    (0, "Grounded, cited, tested Python Q&A — built for trust, not just fluency."),
    (0, "Every answer traceable to a community-verified Stack Overflow source."),
    (0, "Stack: FastAPI · ChromaDB · ONNX MiniLM · DuckDB · Groq Llama 3.1 · pytest."),
    (0, "Repo: github.com/Varuno8/analytics_vidhya_assignment"),
    (0, "Live demo: varun2808/python-qa-assistant (Hugging Face Spaces)"),
], color=RGBColor(0xCB, 0xD5, 0xE1), size=21)

prs.save("slides/deck.pptx")
print("Wrote slides/deck.pptx with", len(prs.slides.slides if hasattr(prs.slides, 'slides') else prs.slides._sldIdLst), "slides")
